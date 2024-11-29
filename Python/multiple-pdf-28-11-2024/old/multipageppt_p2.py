import os
from jinja2 import Environment, FileSystemLoader
from pptx import Presentation
from pptx.util import Pt, Inches
from bs4 import BeautifulSoup

# Directory for templates
TEMPLATE_DIR = "templates"

# Directory for static files (images, CSS)
STATIC_DIR = "static"

# Set up Jinja2 environment
env = Environment(loader=FileSystemLoader(TEMPLATE_DIR))

# Dynamic data for rendering templates
slides_data = [
    {"title": "Introduction", "content": "Welcome to the presentation.", "image_path": "static/page6_graph.png"},
    {"title": "Analysis", "content": "Here is the data analysis.", "image_path": "static/page6_graph.png"},
    {"title": "Conclusion", "content": "Final thoughts and next steps.", "image_path": "static/page6_graph.png"},
]
# Templates list
templates = [
    "template1.html", 
    "template2.html", 
    "template3.html",
    ]
# Create a PowerPoint presentation
presentation = Presentation()

# Function to parse HTML and add content to a slide
def add_html_to_slide(slide, html_content):
    soup = BeautifulSoup(html_content, "html.parser")

    # Add title
    title_element = soup.find("h1")
    if title_element:
        slide.shapes.title.text = title_element.get_text()

    # Add content text
    content_element = soup.find("p")
    if content_element:
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(3)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = content_element.get_text()

    # Add image if present
    image_element = soup.find("img")
    if image_element:
        image_path = os.path.abspath(image_element["src"])
        if os.path.exists(image_path):
            try:
                left = Inches(1)
                top = Inches(5)
                slide.shapes.add_picture(image_path, left, top, width=Inches(4))
            except Exception as e:
                print(f"Error adding image '{image_path}': {e}")
        else:
            print(f"Image not found: {image_path}")

# Iterate through the slide data, render templates, and create slides
# for i, slide_data in enumerate(slides_data):
#     # Load and render the template
#     # template = env.get_template(f"slide_template.html")
#     template = env.get_template(template_name)
#     rendered_html = template.render(
#         title=slide_data["title"],
#         content=slide_data["content"],
#         image=slide_data["image_path"],
#     )

for i, template_name in enumerate(templates):
    # Load and render the template with dynamic content
    template = env.get_template(template_name)
    rendered_html = template.render(
        title=slides_data[i]['title'],
        content=slides_data[i]['content'],
        image_path=slides_data[i]['image_path']
    )

    # Add a new slide
    slide_layout = presentation.slide_layouts[5]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)

    # Populate the slide with rendered HTML content
    add_html_to_slide(slide, rendered_html)

# Save the presentation
output_file = "presentation.pptx"
presentation.save(output_file)
print(f"Presentation saved as {output_file}")