from pptx import Presentation
from pptx.util import Inches, Pt
from jinja2 import Environment, FileSystemLoader
import matplotlib.pyplot as plt
import numpy as np
import os

# Set up Jinja2 environment
template_dir = "templates"
print("Available templates:", os.listdir(template_dir))
env = Environment(loader=FileSystemLoader(template_dir))

# Create a PowerPoint presentation
presentation = Presentation()

# Function to add slide
def add_slide(prs, title, content, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

    # Add title
    if title:
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)

    # Add content
    if content:
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        content_frame = content_box.text_frame
        content_frame.text = content
        content_frame.paragraphs[0].font.size = Pt(18)

    # Add image
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(4), Inches(6), Inches(3))  # Adjust size/position as needed

# Generate the plots
os.makedirs("static", exist_ok=True)
plot_filenames = []

# Page 6 graph
fig, ax = plt.subplots()
x = np.array(["Texto barra 1", "Texto barra 2", "Texto barra 0003"])
y = np.array([2, 4, 7])
ax.bar(x, y, color='#D4C7F3')
plot_filename = "static/page6_graph.png"
plt.savefig(plot_filename, transparent=True)
plt.close()
plot_filenames.append(plot_filename)

# Page 9 graph
sizes = [30, 40, 30]
colors = ['#F36D64', '#FEC84B', '#41CE8C']
plt.pie(sizes, colors=colors, startangle=90, wedgeprops={'width': 0.9})
centre_circle = plt.Circle((0, 0), 0.5, fc='white')
plt.gca().add_artist(centre_circle)
plot_filename = "static/page9_graph.png"
plt.savefig(plot_filename, transparent=True)
plt.close()
plot_filenames.append(plot_filename)

# Dynamic data for slides
dynamic_data = [
    {"title": "Informe de resultados", "content": "Encuesta de clima laboral 2023", "image_path": "static/page6_graph.png"},
    {"title": "Índice", "content": "Contents and agenda of the presentation", "image_path": "static/tamplate09-graph.png"},
    {"title": "Bar Chart Example", "content": "This slide contains a bar chart.", "image_path": "static/tamplate09-graph.png"},
    {"title": "Donut Chart Example", "content": "This slide contains a donut chart.", "image_path": "static/tamplate09-graph.png"},
]

# Templates (optional: matching slide layouts to templates)
templates = ["template1.html", "template2.html", "template3.html", "template4.html"]

# Render HTML and create slides
for i, slide_data in enumerate(dynamic_data):
    # Render template with dynamic content
    if i < len(templates):  # Use the template if it exists
        template = env.get_template(templates[i])
        rendered_html = template.render(
            title=slide_data["title"],
            content=slide_data["content"],
            image_path=slide_data["image_path"],
        )

    # Add slide to PowerPoint
    add_slide(presentation, slide_data["title"], slide_data["content"], slide_data["image_path"])

# Save the presentation
output_pptx_path = "dynamic_presentation.pptx"
presentation.save(output_pptx_path)

print(f"PowerPoint presentation created successfully at {output_pptx_path}!")
