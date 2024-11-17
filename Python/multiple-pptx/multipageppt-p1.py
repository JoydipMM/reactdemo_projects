import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import cairosvg

def convert_svg_to_png(svg_path, png_path):
    """Convert SVG to PNG using cairosvg"""
    cairosvg.svg2png(url=svg_path, write_to=png_path)

def check_and_convert_image(image_path):
    """Ensure image is in a valid format (e.g., PNG, JPEG) and in RGB format"""
    try:
        with Image.open(image_path) as img:
            img = img.convert('RGB')  # Ensure it's in RGB
            img.save(image_path)  # Re-save it to ensure proper format
            print(f"Image saved as {image_path}")
    except Exception as e:
        print(f"Error processing image {image_path}: {e}")

def create_pptx():
    presentation = Presentation()

    dynamic_data = [
        {"title": "Informe de resultados", "content": "Encuesta de clima laboral 2023", "image_path": "static/images/logo-big.svg"},
        {"title": "Índice", "content": "", "image_path": ""},
        {"title": "Welcome to Page 3", "content": "This is the content of the third page.", "image_path": "static/plot3.png"},
        {"title": "Welcome to Page 6", "content": "This is the content of the third page.", "image_path": "static/tamplate06-graph.png"},
        {"title": "Welcome to Page 9", "content": "This is the content of the third page.", "image_path": "static/tamplate09-graph.png"},
    ]

    for data in dynamic_data:
        image_path = data['image_path']

        # Check if the image is an SVG and convert to PNG if needed
        if image_path.endswith(".svg"):
            png_image_path = image_path.replace(".svg", ".png")
            convert_svg_to_png(image_path, png_image_path)
            image_path = png_image_path  # Update image path to PNG version

        # Check and convert image to a valid format
        if os.path.exists(image_path):
            check_and_convert_image(image_path)

            slide_layout = presentation.slide_layouts[5]  # Blank slide
            slide = presentation.slides.add_slide(slide_layout)

            # Title
            title = slide.shapes.title
            title.text = data['title']

            # Content box
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(3)
            content_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = content_box.text_frame
            p = text_frame.add_paragraph()
            p.text = data['content']

            # Add image if it exists
            slide.shapes.add_picture(image_path, Inches(1), Inches(4), width=Inches(8))
        else:
            print(f"Error: Image not found at {image_path}")

    pptx_filename = "generated_presentation.pptx"
    presentation.save(pptx_filename)
    print(f"PowerPoint presentation saved as {pptx_filename}")

# Create the PPTX presentation
create_pptx()